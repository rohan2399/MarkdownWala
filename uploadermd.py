import shutil
import os
import datetime
import pytesseract
from PIL import Image
import fitz
from langchain.schema import Document
import pandas as pd
from dotenv import load_dotenv

# --- Load Environment Variables ---
load_dotenv()

# Langchain Loader
from langchain_community.document_loaders import TextLoader

# Set your Tesseract executable path using the .env file
pytesseract.pytesseract.tesseract_cmd = os.getenv("TESSERACT_CMD", r"C:\Program Files\Tesseract-OCR\tesseract.exe")

# --- Dynamic Folder Configuration ---
# 1. Get the absolute path of the directory where THIS script lives
BASE_DIR = os.path.dirname(os.path.abspath(__file__))

# 2. Fetch the relative paths from the .env file
rel_input = os.getenv("INPUT_FOLDER", "NormalFolder")
rel_obsidian = os.getenv("OBSIDIAN_FOLDER", "ObsidianVault/ExtractedDocs")
rel_attachments = os.getenv("OBSIDIAN_ATTACHMENTS", "ObsidianVault/ExtractedDocs/Attachments")

# 3. Resolve them cleanly to absolute paths on the user's machine
INPUT_FOLDER = os.path.normpath(os.path.join(BASE_DIR, rel_input))
OBSIDIAN_FOLDER = os.path.normpath(os.path.join(BASE_DIR, rel_obsidian))
OBSIDIAN_ATTACHMENTS = os.path.normpath(os.path.join(BASE_DIR, rel_attachments))

# Create all necessary folders automatically if they don't exist yet
os.makedirs(INPUT_FOLDER, exist_ok=True)
os.makedirs(OBSIDIAN_FOLDER, exist_ok=True)
os.makedirs(OBSIDIAN_ATTACHMENTS, exist_ok=True)

print(f"Targeting Input Directory: {INPUT_FOLDER}")
print(f"Targeting Obsidian Directory: {OBSIDIAN_FOLDER}")

# ... [The rest of your script continues exactly as before] ...


def extract_text_from_file(file_path):
    ext = file_path.lower()
    
    if ext.endswith(".txt"):
        # Native Python reading is much safer for handling weird Windows encodings
        try:
            with open(file_path, "r", encoding="utf-8") as f:
                return f.read()
        except UnicodeDecodeError:
            # If UTF-8 fails, fallback to the standard Windows encoding
            with open(file_path, "r", encoding="cp1252", errors="ignore") as f:
                return f.read()
        
    elif ext.endswith(".pdf"):
        # Hybrid Engine: PyMuPDF + pdfplumber + Tesseract OCR Fallback
        import pdfplumber
        
        pdf_document = fitz.open(file_path)
        extracted_content = []
        original_name = os.path.basename(file_path).replace(".pdf", "")
        
        with pdfplumber.open(file_path) as plumb_pdf:
            for page_num in range(len(pdf_document)):
                page = pdf_document.load_page(page_num)
                plumb_page = plumb_pdf.pages[page_num]
                
                # --- 1. TABLE EXTRACTION (Using pdfplumber) ---
                custom_table_settings = {
                    "vertical_strategy": "lines", 
                    "horizontal_strategy": "lines",
                    "intersection_tolerance": 15,
                    "snap_tolerance": 5
                }
                
                tables = plumb_page.find_tables(custom_table_settings)
                table_bboxes = [t.bbox for t in tables] if tables else []
                tables_processed = []
                
                blocks = page.get_text("blocks")
                blocks.sort(key=lambda b: (b[1], b[0]))
                
                page_text_extracted = ""
                
                for block in blocks:
                    if block[6] == 0:
                        bx0, by0, bx1, by1 = block[:4]
                        is_in_table = False
                        
                        for i, t_bbox in enumerate(table_bboxes):
                            tx0, ty0, tx1, ty1 = t_bbox
                            
                            if not (bx1 < tx0 or bx0 > tx1 or by1 < ty0 or by0 > ty1):
                                is_in_table = True
                                
                                if i not in tables_processed:
                                    try:
                                        table_data = tables[i].extract()
                                        if table_data:
                                            df = pd.DataFrame(table_data)
                                            df = df.astype(str)
                                            for col in df.columns:
                                                df[col] = df[col].str.replace('\n', '<br>', regex=False)
                                                df[col] = df[col].str.replace('\r', '', regex=False)
                                                df[col] = df[col].str.replace('None', '', regex=False)
                                            
                                            df.replace('', pd.NA, inplace=True)
                                            df.dropna(axis=0, how='all', inplace=True)
                                            df.dropna(axis=1, how='all', inplace=True)
                                            df.fillna("", inplace=True)
                                            
                                            if not df.empty:
                                                headers = df.iloc[0].values.tolist()
                                                clean_headers = [str(h).strip() if str(h).strip() not in ["", "<br>"] else f"Col_{j+1}" for j, h in enumerate(headers)]
                                                df.columns = clean_headers
                                                df = df.iloc[1:] 
                                                
                                                extracted_table = f"\n{df.to_markdown(index=False)}\n"
                                                extracted_content.append(extracted_table)
                                                page_text_extracted += extracted_table
                                    except Exception as e:
                                        print(f"Error parsing table on page {page_num}: {e}")
                                        
                                    tables_processed.append(i)
                                break 
                        
                        if not is_in_table:
                            text_chunk = block[4]
                            extracted_content.append(text_chunk)
                            page_text_extracted += text_chunk
                
                # --- 2. OCR FALLBACK FOR SCANNED PDFs ---
                if len(page_text_extracted.strip()) < 20:
                    print(f"  -> Page {page_num + 1} of '{original_name}' appears to be image-only. Running OCR...")
                    zoom = 2
                    mat = fitz.Matrix(zoom, zoom)
                    pix = page.get_pixmap(matrix=mat)
                    img = Image.frombytes("RGB", [pix.width, pix.height], pix.samples)
                    ocr_text = pytesseract.image_to_string(img)
                    extracted_content.append(f"\n{ocr_text}\n")

                # --- 3. IMAGE EXTRACTION (Using PyMuPDF) ---
                image_list = page.get_images(full=True)
                for img_index, img in enumerate(image_list):
                    xref = img[0]
                    base_image = pdf_document.extract_image(xref)
                    image_bytes = base_image["image"]
                    image_ext = base_image["ext"]
                    
                    img_filename = f"{original_name}_p{page_num+1}_img{img_index+1}.{image_ext}"
                    img_filepath = os.path.join(OBSIDIAN_ATTACHMENTS, img_filename)
                    
                    with open(img_filepath, "wb") as f:
                        f.write(image_bytes)
                        
                    extracted_content.append(f"\n![[{img_filename}]]\n")
                    
        return "\n".join(extracted_content)

    elif ext.endswith((".docx", ".doc")):
        # On-the-fly PDF Conversion for accurate Word Diagram Extraction
        try:
            from docx2pdf import convert
            
            original_name = os.path.splitext(os.path.basename(file_path))[0]
            print(f"  -> Flattening '{original_name}' to PDF to capture SmartArt/diagrams...")
            
            temp_pdf_path = os.path.join(os.path.dirname(file_path), f"{original_name}_temp_flat.pdf")
            convert(file_path, temp_pdf_path)
            
            extracted_text = extract_text_from_file(temp_pdf_path)
            
            if os.path.exists(temp_pdf_path):
                os.remove(temp_pdf_path)
                
            return extracted_text
            
        except Exception as e:
            print(f"  -> PDF conversion failed (ensure MS Word is closed). Error: {e}")
            return TextLoader(file_path).load()[0].page_content 
        
    elif ext.endswith((".pptx", ".ppt")):
        # --- NEW: NATIVE SLIDE-BY-SLIDE POWERPOINT EXTRACTOR ---
        try:
            from pptx import Presentation
            
            original_name = os.path.splitext(os.path.basename(file_path))[0]
            prs = Presentation(file_path)
            extracted_content = []
            
            for slide_num, slide in enumerate(prs.slides, start=1):
                extracted_content.append(f"\n---\n### Slide {slide_num}\n")
                
                # Recursive function to process text, tables, and images (even if grouped)
                def process_shape(shape):
                    shape_data = []
                    
                    # 1. Standard Text
                    if hasattr(shape, "text") and shape.text.strip():
                        shape_data.append(shape.text.strip())
                        
                    # 2. Native PPT Tables
                    if shape.has_table:
                        table = shape.table
                        table_data = []
                        for row in table.rows:
                            row_data = [cell.text_frame.text.replace('\n', '<br>').replace('\r', '').strip() for cell in row.cells]
                            table_data.append(row_data)
                            
                        if table_data:
                            df = pd.DataFrame(table_data)
                            headers = df.iloc[0].values.tolist()
                            clean_headers = [str(h) if str(h) else f"Col_{j+1}" for j, h in enumerate(headers)]
                            df.columns = clean_headers
                            df = df.iloc[1:]
                            shape_data.append(f"\n{df.to_markdown(index=False)}\n")
                            
                    # 3. Native PPT Images
                    if hasattr(shape, "shape_type") and shape.shape_type == 13: # 13 is MSO_SHAPE_TYPE.PICTURE
                        try:
                            image = shape.image
                            image_bytes = image.blob
                            img_ext = image.ext
                            # Create a clean file name containing the slide number!
                            img_filename = f"{original_name}_Slide{slide_num}_{image.sha1[:8]}.{img_ext}"
                            img_filepath = os.path.join(OBSIDIAN_ATTACHMENTS, img_filename)
                            
                            with open(img_filepath, "wb") as f:
                                f.write(image_bytes)
                                
                            shape_data.append(f"\n![[{img_filename}]]\n")
                        except Exception as e:
                            print(f"Failed to extract image on Slide {slide_num}: {e}")
                            
                    # 4. Grouped Shapes (Dig into the group to find hidden text/images)
                    if hasattr(shape, "shape_type") and shape.shape_type == 6: # 6 is MSO_SHAPE_TYPE.GROUP
                        for sub_shape in shape.shapes:
                            shape_data.extend(process_shape(sub_shape))
                            
                    return shape_data

                for shape in slide.shapes:
                    extracted_content.extend(process_shape(shape))
                    
            return "\n".join(extracted_content)
            
        except Exception as e:
            print(f"Failed to parse PPTX natively, falling back to basic text: {e}")
            return TextLoader(file_path).load()[0].page_content
        
    elif ext.endswith((".xlsx", ".xls")):
        try:
            xls = pd.read_excel(file_path, sheet_name=None)
            extracted_content = []
            for sheet_name, df in xls.items():
                extracted_content.append(f"### Sheet: {sheet_name}\n")
                df.dropna(how='all', inplace=True)
                df.dropna(axis=1, how='all', inplace=True)
                extracted_content.append(df.to_markdown(index=False))
                extracted_content.append("\n")
            return "\n".join(extracted_content)
        except Exception as e:
            print(f"Failed to parse Excel table: {e}")
            return ""       
        
    elif ext.endswith((".jpg", ".jpeg", ".png")):
        # 1. Copy the image to the Obsidian Attachments folder
        img_filename = os.path.basename(file_path)
        img_filepath = os.path.join(OBSIDIAN_ATTACHMENTS, img_filename)
        
        try:
            shutil.copy2(file_path, img_filepath)
        except Exception as e:
            print(f"Failed to copy image {img_filename}: {e}")
            
        # 2. Embed the image in Markdown format
        extracted_content = f"\n![[{img_filename}]]\n\n"
        
        # 3. Run OCR to see if there is any searchable text inside the image
        try:
            ocr_text = pytesseract.image_to_string(Image.open(file_path))
            if ocr_text.strip():
                extracted_content += f"**Image Text:**\n{ocr_text}\n"
        except Exception as e:
            print(f"OCR failed on {img_filename}: {e}")
            
        # Returning this ensures text.strip() evaluates to True because of the Markdown link
        return extracted_content


def extract_all_documents(directory, output_dir):
    documents = []
    for root, _, files in os.walk(directory):
        for f in files:
            file_path = os.path.join(root, f)
            
            if file_path.lower().endswith((".txt", ".pdf", ".docx", ".doc", ".pptx", ".ppt", ".xlsx", ".xls", ".jpg", ".jpeg", ".png")):
                
                rel_dir = os.path.relpath(root, directory)
                target_dir = os.path.join(output_dir, rel_dir)
                os.makedirs(target_dir, exist_ok=True)
                
                base_name = os.path.splitext(f)[0]
                expected_md_path = os.path.join(target_dir, f"{base_name}.md")
                
                if os.path.exists(expected_md_path):
                    source_mtime = os.path.getmtime(file_path)
                    md_mtime = os.path.getmtime(expected_md_path)
                    
                    if source_mtime <= md_mtime:
                        print(f"Skipping: '{f}' (Already up to date in {rel_dir})")
                        continue
                    else:
                        print(f"Updating: '{f}' (Source file was modified)")
                else:
                    print(f"New File Found: '{f}' (Saving to {rel_dir})")

                try:
                    text = extract_text_from_file(file_path)
                    if text.strip():
                        documents.append(Document(page_content=text, metadata={"source": f, "rel_dir": rel_dir}))
                except Exception as e:
                    print(f"Error processing {f}: {e}")
                    
    return documents


def save_to_obsidian(documents, output_dir):
    # This marker separates your personal notes from the script's output.
    # It tells the script where to stop saving and where to start overwriting.
    BOUNDARY = "## 🤖 Auto-Extracted Content"

    for doc in documents:
        original_filename = doc.metadata.get("source", "Unknown_File")
        rel_dir = doc.metadata.get("rel_dir", ".")
        
        base_name = os.path.splitext(original_filename)[0]
        md_filename = f"{base_name}.md"
        
        final_dir = os.path.join(output_dir, rel_dir)
        os.makedirs(final_dir, exist_ok=True)
        
        file_path = os.path.join(final_dir, md_filename)
        
        # Add a timestamp to know when the document was last synced
        current_date = datetime.datetime.now().strftime("%Y-%m-%d %H:%M:%S")
        
        # --- 1. DEFAULT STATE (For brand new files) ---
        user_notes = f"---\noriginal_source: \"{original_filename}\"\ntags: [auto-extracted]\n---\n\n## 📝 My Notes\n*Type your personal notes here...*\n"
        
        # --- 2. PRESERVE EXISTING NOTES (For files being updated) ---
        if os.path.exists(file_path):
            with open(file_path, "r", encoding="utf-8") as md_file:
                existing_content = md_file.read()
            
            # Split the file at the boundary marker
            if BOUNDARY in existing_content:
                # Keep everything ABOVE the marker (Your manual notes!)
                user_notes = existing_content.split(BOUNDARY)[0].strip()
            else:
                # If the marker is missing for some reason, assume the whole 
                # file is your notes so we don't accidentally delete anything.
                user_notes = existing_content.strip()
        
        # --- 3. COMBINE AND SAVE ---
        # We stitch your preserved notes back together with the fresh extraction
        new_file_content = f"{user_notes}\n\n---\n{BOUNDARY}\n**Last Synced:** {current_date}\n\n{doc.page_content}"
        
        # Write it all back to Obsidian
        with open(file_path, "w", encoding="utf-8") as md_file:
            md_file.write(new_file_content)
            
        display_path = os.path.normpath(os.path.join(rel_dir, md_filename))
        print(f"Saved to Obsidian (Notes Preserved): {display_path}")


if __name__ == "__main__":
    print(f"Scanning documents in {INPUT_FOLDER}...")
    extracted_docs = extract_all_documents(INPUT_FOLDER, OBSIDIAN_FOLDER)
    
    if len(extracted_docs) > 0:
        print(f"Found {len(extracted_docs)} new/updated documents. Saving to Obsidian...")
        save_to_obsidian(extracted_docs, OBSIDIAN_FOLDER)
    else:
        print("No new or updated files found.")
        
    print("Done!")
